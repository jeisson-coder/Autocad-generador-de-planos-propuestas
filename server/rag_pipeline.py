"""
rag_pipeline.py — Pipeline RAG para el Mundo Semántico
Procesa documentos del Know How, genera chunks con embeddings (Gemini)
y los almacena en Supabase (pgvector) para búsqueda semántica.

Ejecutar: python rag_pipeline.py
"""

import os
import re
import json
import time
import psycopg2
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from dotenv import load_dotenv
import google.generativeai as genai

load_dotenv()

DATA_ROOT = os.getenv("DATA_ROOT")
KNOWHOW_DIR = os.path.join(DATA_ROOT, "Know How para la IA")
CATALOGOS_DIR = os.path.join(DATA_ROOT, "Catálogos")
GEMINI_API_KEYS = [
    os.getenv("GEMINI_API_KEY"),
    os.getenv("GEMINI_API_KEY_2"),
    os.getenv("GEMINI_API_KEY_3"),
    os.getenv("GEMINI_API_KEY_4"),
    os.getenv("GEMINI_API_KEY_5"),
    os.getenv("GEMINI_API_KEY_6"),
]
GEMINI_API_KEYS = [k for k in GEMINI_API_KEYS if k]  # Filtrar None
_current_key_idx = 0

def _configure_gemini(key_idx=None):
    """Configura Gemini con la key actual o la siguiente."""
    global _current_key_idx
    if key_idx is not None:
        _current_key_idx = key_idx
    genai.configure(api_key=GEMINI_API_KEYS[_current_key_idx % len(GEMINI_API_KEYS)])
    return _current_key_idx % len(GEMINI_API_KEYS)

_configure_gemini(0)

# ============================================================
# CATEGORÍAS POR NOMBRE DE ARCHIVO
# ============================================================

CATEGORIA_MAP = {
    "sistema restauración": "diseno_general",
    "autocontrol": "normativa",
    "LIMPIEZA": "limpieza",
    "restaurantes rec": "tipos_negocio",
    "restaurantes": "tipos_negocio",
    "diseño a través de un plano": "diseno_general",
    "Taperia cafeteria": "tipos_negocio",
    "RECEPCIÓN": "zonas_cocina",
    "barras tapas": "tipos_negocio",
    "ALMACENAMIENTO": "zonas_cocina",
    "ALMACENAMENTO FRIGORÍFICO": "zonas_cocina",
    "pizzerias": "tipos_negocio",
    "MANTENIMIENTO CAMARAS": "zonas_cocina",
    "fast food": "tipos_negocio",
    "DESCONGELACIÓN": "zonas_cocina",
    "CUARTO FRIO": "zonas_cocina",
    "orientales": "tipos_negocio",
    "CUARTO FRIO ELABORADOS": "zonas_cocina",
    "ZONA DE COCCIÓN": "zonas_cocina",
    "ENFRIAMIENTO": "zonas_cocina",
    "LAVADO DE VAJILLA": "zonas_cocina",
    "PLONGE": "zonas_cocina",
    "CUARTO LIMPIEZA": "zonas_cocina",
    "CUARTO DE BASURAS": "zonas_cocina",
    "BUFE": "zonas_cocina",
    "SHOW COOKING": "zonas_cocina",
    "paredes y tabiques": "construccion",
    "SUELOS Y DESAGUES": "construccion",
    "VENTILACIÓN": "construccion",
    "aplicaciones gastronómicas": "hornos",
    "SALONES DE BANQUETES": "tipos_negocio",
    "CATERING": "tipos_negocio",
    "checklist toma de datos": "diseno_general",
    "dimensiones-reducidas": "diseno_general",
    "diseño-cocinas-hospitales": "tipos_negocio",
    "legislacion": "normativa",
    "LEGISLACIÓN": "normativa",
    "dimensionado": "dimensionado",
    "COMENSALES": "dimensionado",
    "RATIONAL": "catalogo",
    "Serie": "catalogo",
    "StarLine": "catalogo",
    "Accesorios": "catalogo",
    "Varios": "catalogo",
    "Apice": "catalogo",
    "Retigo": "catalogo",
    "retigo": "catalogo",
    "Nikrom": "catalogo",
    "SARA": "catalogo",
    "GEMM": "catalogo",
    "manualmarca": "identidad_marca",
    "manual de marca": "identidad_marca",
    "PROSPECCIÓN": "prospeccion",
    "prospeccion": "prospeccion",
}


def detectar_categoria(filename: str) -> str:
    """Detecta la categoría del documento por su nombre."""
    for keyword, cat in CATEGORIA_MAP.items():
        if keyword.lower() in filename.lower():
            return cat
    return "general"


# ============================================================
# EXTRACTORES DE TEXTO
# ============================================================

def extraer_texto_pdf(path: str) -> str:
    """Extrae texto de un PDF. Retorna string vacío si es escaneado."""
    try:
        reader = PdfReader(path)
        textos = []
        for page in reader.pages:
            text = page.extract_text()
            if text and text.strip():
                textos.append(text.strip())
        return "\n\n".join(textos)
    except Exception as e:
        print(f"    WARN: Error leyendo PDF {os.path.basename(path)}: {e}")
        return ""


def extraer_texto_pptx(path: str) -> str:
    """Extrae texto de un PPTX/PPTM."""
    try:
        prs = Presentation(path)
        textos = []
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
            if slide_texts:
                textos.append(f"[Slide {i+1}]\n" + "\n".join(slide_texts))
        return "\n\n".join(textos)
    except Exception as e:
        print(f"    WARN: Error leyendo PPTX {os.path.basename(path)}: {e}")
        return ""


def extraer_texto_docx(path: str) -> str:
    """Extrae texto de un DOCX."""
    try:
        doc = Document(path)
        textos = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                textos.append(text)
        return "\n\n".join(textos)
    except Exception as e:
        print(f"    WARN: Error leyendo DOCX {os.path.basename(path)}: {e}")
        return ""


# ============================================================
# CHUNKING
# ============================================================

def chunkear_texto(texto: str, chunk_size: int = 800, overlap: int = 200) -> list[str]:
    """Divide texto en chunks con overlap."""
    if not texto or len(texto.strip()) < 50:
        return []

    # Limpiar texto
    texto = re.sub(r"\n{3,}", "\n\n", texto)
    texto = re.sub(r" {2,}", " ", texto)

    # Dividir por dobles saltos de línea (párrafos)
    parrafos = texto.split("\n\n")

    chunks = []
    chunk_actual = ""

    for parrafo in parrafos:
        parrafo = parrafo.strip()
        if not parrafo:
            continue

        # Si agregar este párrafo excede el tamaño, cerrar chunk actual
        if len(chunk_actual) + len(parrafo) > chunk_size and chunk_actual:
            chunks.append(chunk_actual.strip())
            # Overlap: tomar últimos N caracteres del chunk anterior
            if overlap > 0 and len(chunk_actual) > overlap:
                chunk_actual = chunk_actual[-overlap:] + "\n\n" + parrafo
            else:
                chunk_actual = parrafo
        else:
            if chunk_actual:
                chunk_actual += "\n\n" + parrafo
            else:
                chunk_actual = parrafo

    # Último chunk
    if chunk_actual.strip():
        chunks.append(chunk_actual.strip())

    # Subdividir chunks que siguen siendo muy grandes
    final_chunks = []
    for chunk in chunks:
        if len(chunk) > chunk_size * 2:
            # Dividir por oraciones
            sentences = re.split(r"(?<=[.!?])\s+", chunk)
            sub_chunk = ""
            for sent in sentences:
                if len(sub_chunk) + len(sent) > chunk_size and sub_chunk:
                    final_chunks.append(sub_chunk.strip())
                    sub_chunk = sent
                else:
                    sub_chunk += " " + sent if sub_chunk else sent
            if sub_chunk.strip():
                final_chunks.append(sub_chunk.strip())
        else:
            final_chunks.append(chunk)

    return final_chunks


# ============================================================
# EMBEDDINGS CON GEMINI
# ============================================================

def generar_embeddings(textos: list[str], batch_size: int = 10) -> list[list[float]]:
    """Genera embeddings con Gemini gemini-embedding-001.
    Rota entre 4 API keys proactivamente en cada batch.
    Batch size=10 con 4 keys = ~400 req/min capacidad."""
    global _current_key_idx
    all_embeddings = []
    total_batches = (len(textos) + batch_size - 1) // batch_size

    for i in range(0, len(textos), batch_size):
        batch = textos[i : i + batch_size]
        batch_num = i // batch_size + 1

        # Rotar keys proactivamente en cada batch
        _current_key_idx = (batch_num - 1) % len(GEMINI_API_KEYS)
        _configure_gemini(_current_key_idx)

        retries = 0
        max_retries = 20

        while retries < max_retries:
            try:
                result = genai.embed_content(
                    model="models/gemini-embedding-001",
                    content=batch,
                    output_dimensionality=768,
                )
                all_embeddings.extend(result["embedding"])
                if batch_num % 10 == 0 or batch_num == total_batches:
                    print(f"      Batch {batch_num}/{total_batches} OK [key {_current_key_idx}]", flush=True)
                break
            except Exception as e:
                retries += 1
                error_str = str(e)

                if "429" in error_str:
                    # Rotar a la otra key
                    _current_key_idx = (_current_key_idx + 1) % len(GEMINI_API_KEYS)
                    _configure_gemini(_current_key_idx)

                    # Extraer tiempo de espera
                    import re as _re
                    delay_match = _re.search(r"retry in (\d+)", error_str)
                    wait = int(delay_match.group(1)) + 5 if delay_match else 35
                    wait = min(wait, 65)  # Limitar espera maxima a 65s
                    print(f"      Batch {batch_num} rate limited, esperando {wait}s... [key {_current_key_idx}]", flush=True)
                    time.sleep(wait)
                else:
                    wait = min(2 ** retries, 60)
                    print(f"      Batch {batch_num} error, retry {retries}/{max_retries} en {wait}s: {e}", flush=True)
                    time.sleep(wait)

                if retries >= max_retries:
                    print(f"      ERROR: Batch {batch_num} falló despues de {max_retries} intentos", flush=True)
                    all_embeddings.extend([[0.0] * 768] * len(batch))

        # Pausa entre batches: 3s para ~100 req/min con 2 keys
        if i + batch_size < len(textos):
            time.sleep(3)

    return all_embeddings


# ============================================================
# ALMACENAMIENTO EN SUPABASE
# ============================================================

def get_db_connection():
    """Crea conexión a Supabase PostgreSQL."""
    return psycopg2.connect(
        host=os.getenv("SUPABASE_DB_HOST"),
        port=os.getenv("SUPABASE_DB_PORT"),
        dbname=os.getenv("SUPABASE_DB_NAME"),
        user=os.getenv("SUPABASE_DB_USER"),
        password=os.getenv("SUPABASE_DB_PASSWORD"),
        sslmode="require",
    )


def almacenar_documento(cur, titulo, tipo_archivo, ruta, categoria, chunks, embeddings):
    """Almacena un documento y sus chunks en Supabase."""
    # Insertar documento
    cur.execute(
        """INSERT INTO documentos_rag (titulo, tipo_archivo, ruta_origen, categoria, num_chunks, procesado)
           VALUES (%s, %s, %s, %s, %s, TRUE)
           RETURNING id;""",
        (titulo, tipo_archivo, ruta, categoria, len(chunks)),
    )
    doc_id = cur.fetchone()[0]

    # Insertar chunks con embeddings
    for idx, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
        metadata = json.dumps({"chunk_index": idx, "categoria": categoria})
        cur.execute(
            """INSERT INTO chunks_rag (documento_id, chunk_index, contenido, embedding, metadata)
               VALUES (%s, %s, %s, %s::vector, %s);""",
            (str(doc_id), idx, chunk, str(embedding), metadata),
        )

    return doc_id


# ============================================================
# BÚSQUEDA SEMÁNTICA
# ============================================================

def buscar_similar(query: str, top_k: int = 5, categoria: str = None) -> list[dict]:
    """Busca chunks similares a la query usando embeddings Gemini."""
    # Generar embedding de la query
    result = genai.embed_content(
        model="models/gemini-embedding-001",
        content=query,
        output_dimensionality=768,
    )
    query_embedding = result["embedding"]

    conn = get_db_connection()
    cur = conn.cursor()

    if categoria:
        cur.execute(
            """SELECT c.contenido, d.titulo, d.categoria,
                      1 - (c.embedding <=> %s::vector) as similitud
               FROM chunks_rag c
               JOIN documentos_rag d ON c.documento_id = d.id
               WHERE d.categoria = %s
               ORDER BY c.embedding <=> %s::vector
               LIMIT %s;""",
            (str(query_embedding), categoria, str(query_embedding), top_k),
        )
    else:
        cur.execute(
            """SELECT c.contenido, d.titulo, d.categoria,
                      1 - (c.embedding <=> %s::vector) as similitud
               FROM chunks_rag c
               JOIN documentos_rag d ON c.documento_id = d.id
               ORDER BY c.embedding <=> %s::vector
               LIMIT %s;""",
            (str(query_embedding), str(query_embedding), top_k),
        )

    results = []
    for row in cur.fetchall():
        results.append({
            "contenido": row[0],
            "titulo": row[1],
            "categoria": row[2],
            "similitud": float(row[3]),
        })

    cur.close()
    conn.close()
    return results


# ============================================================
# MAIN
# ============================================================

def main():
    print("=" * 60)
    print("RAG Pipeline — Mundo Semántico Repagas")
    print("=" * 60)

    # Recolectar archivos procesables
    archivos = []

    # Know How para la IA
    print(f"\n--- Escaneando {KNOWHOW_DIR} ---")
    if os.path.exists(KNOWHOW_DIR):
        for f in os.listdir(KNOWHOW_DIR):
            full = os.path.join(KNOWHOW_DIR, f)
            if not os.path.isfile(full):
                continue
            ext = os.path.splitext(f)[1].lower()
            if ext in (".pdf", ".pptx", ".pptm", ".docx"):
                archivos.append((full, ext, "knowhow"))
                print(f"  + {f} ({ext})")
            elif ext in (".ppt", ".doc"):
                print(f"  - {f} ({ext}) [SKIP: formato antiguo, necesita LibreOffice]")

    # Logo / Manual de Marca + PPTX Prospección
    MARCA_DIR = os.path.join(DATA_ROOT, "Logo-Manual de Marca")
    print(f"\n--- Escaneando {MARCA_DIR} ---")
    if os.path.exists(MARCA_DIR):
        for f in os.listdir(MARCA_DIR):
            full = os.path.join(MARCA_DIR, f)
            if not os.path.isfile(full):
                continue
            ext = os.path.splitext(f)[1].lower()
            if ext in (".pdf", ".pptx", ".pptm", ".docx"):
                archivos.append((full, ext, "marca"))
                print(f"  + {f} ({ext})")

    # PPTX Prospección (archivo suelto en la raíz de Repagas IA)
    prospeccion_pptx = os.path.join(DATA_ROOT, "REPAGAS CONCEPT PROSPECCIÓN.pptx")
    if os.path.isfile(prospeccion_pptx):
        archivos.append((prospeccion_pptx, ".pptx", "prospeccion"))
        print(f"  + REPAGAS CONCEPT PROSPECCIÓN.pptx")

    # Catálogos
    print(f"\n--- Escaneando {CATALOGOS_DIR} ---")
    if os.path.exists(CATALOGOS_DIR):
        for f in os.listdir(CATALOGOS_DIR):
            full = os.path.join(CATALOGOS_DIR, f)
            ext = os.path.splitext(f)[1].lower()
            if ext == ".pdf" and os.path.isfile(full):
                archivos.append((full, ext, "catalogo"))
                print(f"  + {f}")

    print(f"\n  Total archivos a procesar: {len(archivos)}")

    # Verificar docs ya procesados
    conn = get_db_connection()
    cur = conn.cursor()
    cur.execute("SELECT ruta_origen FROM documentos_rag WHERE procesado = TRUE")
    ya_procesados = set(row[0] for row in cur.fetchall())
    cur.close()
    conn.close()

    print(f"\n  Documentos ya en DB: {len(ya_procesados)}")
    archivos_nuevos = [(fp, ext, src) for fp, ext, src in archivos if fp not in ya_procesados]
    print(f"  Archivos nuevos a procesar: {len(archivos_nuevos)}")

    if not archivos_nuevos:
        print("\n  No hay archivos nuevos. Todo esta al dia.")

    # Procesar cada archivo
    total_chunks = 0
    docs_procesados = 0
    docs_sin_texto = 0
    stats_por_categoria = {}

    for filepath, ext, source in archivos_nuevos:
        filename = os.path.basename(filepath)
        print(f"\n  [{docs_procesados+1}/{len(archivos_nuevos)}] {filename}")

        # Extraer texto
        if ext == ".pdf":
            texto = extraer_texto_pdf(filepath)
        elif ext in (".pptx", ".pptm"):
            texto = extraer_texto_pptx(filepath)
        elif ext == ".docx":
            texto = extraer_texto_docx(filepath)
        else:
            continue

        if not texto or len(texto.strip()) < 50:
            print(f"    SKIP: sin texto extraible ({len(texto) if texto else 0} chars)")
            docs_sin_texto += 1
            continue

        print(f"    Texto: {len(texto)} chars")

        # Chunkear
        chunks = chunkear_texto(texto)
        if not chunks:
            print(f"    SKIP: 0 chunks generados")
            docs_sin_texto += 1
            continue

        print(f"    Chunks: {len(chunks)}")

        # Generar embeddings (puede tardar por rate limits)
        print(f"    Generando embeddings...")
        embeddings = generar_embeddings(chunks)

        # Detectar categoria
        categoria = detectar_categoria(filename)
        tipo_archivo = ext.replace(".", "")

        # Almacenar con conexion FRESCA (evitar timeout por idle)
        titulo = os.path.splitext(filename)[0]
        try:
            conn = get_db_connection()
            conn.autocommit = False
            cur = conn.cursor()
            almacenar_documento(cur, titulo, tipo_archivo, filepath, categoria, chunks, embeddings)
            conn.commit()
            cur.close()
            conn.close()
        except Exception as e:
            print(f"    ERROR almacenando: {e}")
            try:
                conn.rollback()
                cur.close()
                conn.close()
            except Exception:
                pass
            continue

        total_chunks += len(chunks)
        docs_procesados += 1
        stats_por_categoria[categoria] = stats_por_categoria.get(categoria, 0) + len(chunks)
        print(f"    OK: {len(chunks)} chunks [{categoria}]")

    # Resumen final
    print("\n" + "=" * 60)
    print("RESUMEN RAG PIPELINE")
    print("=" * 60)
    print(f"  Documentos procesados: {docs_procesados}")
    print(f"  Documentos sin texto:  {docs_sin_texto}")
    print(f"  Total chunks:          {total_chunks}")
    print(f"\n  Chunks por categoria:")
    for cat, count in sorted(stats_por_categoria.items()):
        print(f"    {cat}: {count}")

    # Verificacion en DB
    conn = get_db_connection()
    cur = conn.cursor()
    cur.execute("SELECT COUNT(*) FROM documentos_rag WHERE procesado = TRUE;")
    docs_db = cur.fetchone()[0]
    cur.execute("SELECT COUNT(*) FROM chunks_rag;")
    chunks_db = cur.fetchone()[0]
    print(f"\n  Verificacion Supabase:")
    print(f"    documentos_rag: {docs_db}")
    print(f"    chunks_rag: {chunks_db}")
    cur.close()
    conn.close()

    print("\n" + "=" * 60)
    print("Pipeline RAG completado.")
    print("=" * 60)

if __name__ == "__main__":
    main()


